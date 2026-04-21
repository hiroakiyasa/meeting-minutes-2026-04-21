# -*- coding: utf-8 -*-
"""jp-automation-pptx の PPTX ビルダー。
Automation スタイル（日本企業標準）を基盤に、data-viz と mermaid を必要に応じて組み込む。

入力: input.json
出力: output.pptx
"""
import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Same dir imports
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from render_chart import (bar_horizontal, bar_compare, bar_vertical, gantt,
                           priority_matrix)
from render_mermaid import render_mermaid

# ---------- Theme colors ----------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
BLUE = RGBColor(0x2E, 0x6B, 0xC0)
LIGHT = RGBColor(0xE9, 0xEF, 0xF7)
TEXT = RGBColor(0x22, 0x26, 0x2E)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC4, 0x4E, 0x52)


def emu_zero():
    return Emu(0)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.W = self.prs.slide_width
        self.H = self.prs.slide_height

    # ---------- primitives ----------
    def add_header(self, slide, text):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu_zero(), emu_zero(),
                                     self.W, Inches(0.9))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = text; r.font.size = Pt(22); r.font.bold = True
        r.font.name = "游ゴシック"; r.font.color.rgb = WHITE

    def add_text(self, slide, l, t, w, h, text, size, bold=False, color=TEXT,
                 font="游ゴシック", align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = line
            r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color

    def add_rect(self, slide, l, t, w, h, fill, line=None, rounded=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(shape_type, l, t, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line:
            sh.line.color.rgb = line
        else:
            sh.line.fill.background()
        return sh

    def add_image(self, slide, path, l, t, w=None, h=None):
        if w is not None:
            slide.shapes.add_picture(path, l, t, width=w)
        elif h is not None:
            slide.shapes.add_picture(path, l, t, height=h)
        else:
            slide.shapes.add_picture(path, l, t)

    def new(self):
        return self.prs.slides.add_slide(self.blank)

    # ---------- slide templates ----------
    def slide_cover(self, title, subtitle, footer=""):
        s = self.new()
        self.add_rect(s, emu_zero(), emu_zero(), self.W, self.H, LIGHT)
        self.add_rect(s, emu_zero(), Inches(2.5), self.W, Inches(2.6), NAVY)
        self.add_text(s, Inches(0.8), Inches(2.9), Inches(12), Inches(1.2),
                      title, 40, bold=True, color=WHITE)
        self.add_text(s, Inches(0.8), Inches(4.1), Inches(12), Inches(0.6),
                      subtitle, 16, color=RGBColor(0xD5, 0xE1, 0xF5))
        if footer:
            self.add_text(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
                          footer, 10, color=GRAY)

    def slide_summary(self, header, kpis, sub_blocks=None):
        """kpis = [(label, big, sub), ...]  sub_blocks = [(title, body), ...]"""
        s = self.new()
        self.add_header(s, header)
        n = len(kpis)
        total_w = 12.3; margin = 0.5
        card_w = (total_w - (n - 1) * 0.2) / n
        for i, (label, big, sub) in enumerate(kpis):
            x = Inches(margin + i * (card_w + 0.2))
            self.add_rect(s, x, Inches(1.4), Inches(card_w), Inches(2.5),
                          WHITE, line=BLUE, rounded=True)
            self.add_text(s, x + Inches(0.2), Inches(1.55), Inches(card_w - 0.4),
                          Inches(0.4), label, 12, bold=True, color=BLUE)
            self.add_text(s, x + Inches(0.2), Inches(2.0), Inches(card_w - 0.4),
                          Inches(1.0), big, 40, bold=True, color=NAVY)
            self.add_text(s, x + Inches(0.2), Inches(3.2), Inches(card_w - 0.4),
                          Inches(0.5), sub, 11, color=GRAY)
        if sub_blocks:
            for i, (stitle, sbody) in enumerate(sub_blocks):
                bw = (total_w - (len(sub_blocks) - 1) * 0.2) / len(sub_blocks)
                x = Inches(margin + i * (bw + 0.2))
                self.add_rect(s, x, Inches(4.9), Inches(bw), Inches(1.9), LIGHT)
                self.add_text(s, x + Inches(0.2), Inches(5.1), Inches(bw - 0.4),
                              Inches(0.5), stitle, 16, bold=True, color=NAVY)
                self.add_text(s, x + Inches(0.2), Inches(5.7), Inches(bw - 0.4),
                              Inches(1.1), sbody, 12)

    def slide_agenda(self, items):
        s = self.new()
        self.add_header(s, "アジェンダ")
        for i, item in enumerate(items):
            self.add_rect(s, Inches(0.8), Inches(1.5 + i * 1.5), Inches(11.7),
                          Inches(1.2), LIGHT, line=BLUE, rounded=True)
            self.add_text(s, Inches(1.1), Inches(1.75 + i * 1.5), Inches(1.5),
                          Inches(0.8), f"{i+1}", 44, bold=True, color=BLUE)
            self.add_text(s, Inches(2.4), Inches(1.9 + i * 1.5), Inches(10),
                          Inches(0.6), item, 18, bold=True)

    def slide_section_text(self, header, lead, bullets):
        """テキスト主体スライド"""
        s = self.new()
        self.add_header(s, header)
        self.add_text(s, Inches(0.6), Inches(1.2), Inches(12.5), Inches(0.6),
                      "▼ " + lead, 15, bold=True, color=BLUE)
        self.add_rect(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(5.2), LIGHT)
        for i, b in enumerate(bullets):
            self.add_text(s, Inches(0.9), Inches(2.1 + i * 0.75), Inches(11.7),
                          Inches(0.7), "● " + b, 13)

    def slide_section_with_chart(self, header, lead, bullets, chart_path,
                                  chart_caption=""):
        """テキスト + 右にチャート（data-viz連携）"""
        s = self.new()
        self.add_header(s, header)
        self.add_text(s, Inches(0.6), Inches(1.2), Inches(12.5), Inches(0.6),
                      "▼ " + lead, 15, bold=True, color=BLUE)
        # 左側テキスト
        self.add_rect(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.2), LIGHT)
        for i, b in enumerate(bullets):
            self.add_text(s, Inches(0.7), Inches(2.1 + i * 0.7), Inches(5.5),
                          Inches(0.6), "● " + b, 11)
        # 右側チャート
        self.add_image(s, chart_path, Inches(6.6), Inches(1.9), w=Inches(6.5))
        if chart_caption:
            self.add_text(s, Inches(6.6), Inches(6.9), Inches(6.5), Inches(0.4),
                          "▸ " + chart_caption, 10, color=GRAY)

    def slide_section_with_mermaid(self, header, lead, bullets, mermaid_path):
        """テキスト + 右にMermaid図"""
        s = self.new()
        self.add_header(s, header)
        self.add_text(s, Inches(0.6), Inches(1.2), Inches(12.5), Inches(0.6),
                      "▼ " + lead, 15, bold=True, color=BLUE)
        self.add_rect(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.2), LIGHT)
        for i, b in enumerate(bullets):
            self.add_text(s, Inches(0.7), Inches(2.1 + i * 0.7), Inches(5.5),
                          Inches(0.6), "● " + b, 11)
        self.add_image(s, mermaid_path, Inches(6.8), Inches(2.0), w=Inches(6.3))

    def slide_full_chart(self, header, chart_path, note=""):
        s = self.new()
        self.add_header(s, header)
        self.add_image(s, chart_path, Inches(0.8), Inches(1.2), w=Inches(11.7))
        if note:
            self.add_text(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5),
                          "▼ " + note, 12, bold=True, color=NAVY)

    def slide_decisions(self, decisions):
        s = self.new()
        self.add_header(s, "決定事項")
        for i, d in enumerate(decisions):
            self.add_rect(s, Inches(0.6), Inches(1.25 + i * 0.88), Inches(12.1),
                          Inches(0.78),
                          LIGHT if i % 2 == 0 else WHITE, line=BLUE)
            self.add_text(s, Inches(0.8), Inches(1.4 + i * 0.88), Inches(1),
                          Inches(0.5), f"D{i+1}", 18, bold=True, color=BLUE)
            self.add_text(s, Inches(1.9), Inches(1.45 + i * 0.88), Inches(10.8),
                          Inches(0.5), d, 13)

    def slide_todos(self, todos, extra_chart_path=None):
        """todos = [(owner, action, due), ...]"""
        s = self.new()
        self.add_header(s, "ToDo（次回までのアクション）")
        if extra_chart_path:
            # 左チャート・右テーブル
            self.add_image(s, extra_chart_path, Inches(0.4), Inches(1.1), w=Inches(6.3))
            table_l = Inches(7.0)
            table_w = Inches(6.0)
        else:
            table_l = Inches(0.6)
            table_w = Inches(12.1)
        rows = len(todos) + 1
        table = s.shapes.add_table(rows, 3, table_l, Inches(1.3), table_w,
                                    Inches(5.5)).table
        col_ratios = [0.15, 0.65, 0.2]
        for i, r_ in enumerate(col_ratios):
            table.columns[i].width = int(table_w * r_)
        for i, h in enumerate(["担当", "アクション", "期限"]):
            cell = table.cell(0, i); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.size = Pt(11)
                    r.font.color.rgb = WHITE; r.font.name = "游ゴシック"
        for i, (o, a, d) in enumerate(todos, 1):
            for j, v in enumerate([o, a, d]):
                cell = table.cell(i, j); cell.text = v
                if i % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10); r.font.name = "游ゴシック"

    def slide_closing(self, title, subtitle, footer=""):
        s = self.new()
        self.add_rect(s, emu_zero(), emu_zero(), self.W, self.H, NAVY)
        self.add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(1.5),
                      title, 44, bold=True, color=WHITE)
        self.add_text(s, Inches(0.8), Inches(4.1), Inches(12), Inches(0.7),
                      subtitle, 20, color=RGBColor(0xD5, 0xE1, 0xF5))
        if footer:
            self.add_text(s, Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
                          footer, 10, color=RGBColor(0xA3, 0xB8, 0xD8))

    def save(self, path):
        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
        self.prs.save(path)
        return path


def build_from_json(input_path, output_path, work_dir=None):
    """Main entry: input.json → output.pptx"""
    with open(input_path, encoding="utf-8") as f:
        data = json.load(f)
    work_dir = work_dir or os.path.join(os.path.dirname(os.path.abspath(output_path)),
                                         "_charts")
    os.makedirs(work_dir, exist_ok=True)

    deck = Deck()

    # Cover
    cov = data["cover"]
    deck.slide_cover(cov["title"], cov.get("subtitle", ""), cov.get("footer", ""))

    # Summary
    if "summary" in data:
        deck.slide_summary(data["summary"].get("header", "エグゼクティブサマリー"),
                           [tuple(k) for k in data["summary"]["kpis"]],
                           [tuple(b) for b in data["summary"].get("blocks", [])])

    # Agenda
    if "agenda" in data:
        deck.slide_agenda(data["agenda"])

    # Sections
    for sec in data.get("sections", []):
        visual = sec.get("visual")
        if visual is None:
            deck.slide_section_text(sec["header"], sec["lead"], sec["bullets"])
        elif visual["type"] == "bar_horizontal":
            p = bar_horizontal(visual["data"], visual["title"],
                               visual.get("xlabel", ""),
                               os.path.join(work_dir, f"{visual['id']}.png"),
                               visual.get("unit", ""))
            deck.slide_section_with_chart(sec["header"], sec["lead"], sec["bullets"],
                                          p, visual.get("caption", ""))
        elif visual["type"] == "bar_compare":
            p = bar_compare(visual["categories"], visual["current"], visual["target"],
                            visual["title"],
                            os.path.join(work_dir, f"{visual['id']}.png"),
                            visual.get("ylabel", "%"),
                            visual.get("percent", True))
            deck.slide_section_with_chart(sec["header"], sec["lead"], sec["bullets"],
                                          p, visual.get("caption", ""))
        elif visual["type"] == "priority_matrix":
            p = priority_matrix([tuple(i) for i in visual["items"]], visual["title"],
                                 os.path.join(work_dir, f"{visual['id']}.png"))
            deck.slide_section_with_chart(sec["header"], sec["lead"], sec["bullets"],
                                          p, visual.get("caption", ""))
        elif visual["type"] == "mermaid":
            p = render_mermaid(visual["code"],
                                os.path.join(work_dir, f"{visual['id']}.png"))
            deck.slide_section_with_mermaid(sec["header"], sec["lead"],
                                             sec["bullets"], p)

    # Full-slide visuals
    for fv in data.get("full_visuals", []):
        if fv["type"] == "gantt":
            p = gantt([tuple(t) for t in fv["tasks"]], fv["title"],
                       os.path.join(work_dir, f"{fv['id']}.png"),
                       fv.get("marker_date"), fv.get("marker_label", ""))
            deck.slide_full_chart(fv.get("header", fv["title"]), p, fv.get("note", ""))
        elif fv["type"] == "mermaid":
            p = render_mermaid(fv["code"],
                                os.path.join(work_dir, f"{fv['id']}.png"))
            deck.slide_full_chart(fv.get("header", "図解"), p, fv.get("note", ""))
        elif fv["type"] == "priority_matrix":
            p = priority_matrix([tuple(i) for i in fv["items"]], fv["title"],
                                 os.path.join(work_dir, f"{fv['id']}.png"))
            deck.slide_full_chart(fv.get("header", fv["title"]), p, fv.get("note", ""))

    # Decisions
    if "decisions" in data:
        deck.slide_decisions(data["decisions"])

    # ToDo
    if "todos" in data:
        chart_path = None
        if data["todos"].get("with_load_chart"):
            owners = [t[0] for t in data["todos"]["items"]]
            from collections import Counter
            cnt = Counter(owners)
            chart_path = bar_vertical(dict(cnt), "担当別 ToDo件数",
                                       os.path.join(work_dir, "todo_load.png"),
                                       "件数")
        deck.slide_todos([tuple(t) for t in data["todos"]["items"]], chart_path)

    # Closing
    if "closing" in data:
        c = data["closing"]
        deck.slide_closing(c["title"], c.get("subtitle", ""), c.get("footer", ""))

    return deck.save(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("usage: build_pptx.py <input.json> <output.pptx>")
        sys.exit(1)
    out = build_from_json(sys.argv[1], sys.argv[2])
    print(f"[done] {out} ({os.path.getsize(out) // 1024} KB)")
